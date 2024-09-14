from rest_framework import status
from rest_framework.response import Response
from django.shortcuts import get_object_or_404
from .serializers import UploadedFileSerializer
from pptx import Presentation
import logging
logger = logging.getLogger(__name__)
import base64
import os
from .serializers import HtmlContentSerializer
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Inches
from PIL import Image, ImageDraw
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from io import BytesIO
from django.http import HttpResponse
import requests
from rest_framework.decorators import api_view, permission_classes
from rest_framework.permissions import IsAuthenticated, AllowAny
from django.contrib.auth import get_user_model
from .models import UploadedFile, HtmlContent
from rest_framework.permissions import IsAuthenticated
from rest_framework.decorators import api_view, permission_classes
from django.contrib.auth.models import User  # For default user model
# or if using a custom user model:
# from django.contrib.auth import get_user_model
# User = get_user_model()
from django.contrib.auth import get_user_model
from rest_framework.decorators import api_view
from rest_framework.response import Response
from rest_framework import status
from .models import UploadedFile
from .serializers import UploadedFileSerializer
from rest_framework.decorators import api_view
from rest_framework.response import Response
from rest_framework import status
from .models import UploadedFile
from .serializers import UploadedFileSerializer

User = get_user_model() 

def add_hyperlink_box(pptx_file,output_file):
    prs = Presentation(pptx_file)

    def check_shape_for_hyperlink(shape, slide):
        has_link = False
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.hyperlink.address is not None:
                        has_link = True

        if shape.click_action.hyperlink.address or has_link:
            # Get shape position and size
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            # Add a rectangle around the shape
            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left,
                top,
                width,
                height
            )
            # Set rectangle properties
            line = rect.line
            line.color.rgb = RGBColor(255, 0, 0)  # Red color
            line.width = Pt(3)
            rect.fill.background()  # Transparent fill

    def process_shape(shape, slide):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for subshape in shape.shapes:
                process_shape(subshape, slide)
        else:
            check_shape_for_hyperlink(shape, slide)

    for slide in prs.slides:
        for shape in slide.shapes:
            process_shape(shape, slide)

    prs.save(output_file)
    

@api_view(['POST'])
def save_guide(request):
    
    serializer = HtmlContentSerializer(data=request.data)
    if serializer.is_valid():
        serializer.save()
        return Response(serializer.data, status=status.HTTP_201_CREATED)
    return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

@api_view(['GET'])
@permission_classes([AllowAny])
def list_guides(request):
    html_contents = HtmlContent.objects.all()
    serializer = HtmlContentSerializer(html_contents, many=True)
    return Response(serializer.data, status=status.HTTP_200_OK)


@api_view(['GET'])
@permission_classes([AllowAny])
def process_pptx(request, file_id):
    uploaded_file = get_object_or_404(UploadedFile, id=file_id)
    file_path = uploaded_file.file.path

    # Check the file extension and convert if necessary
    file_ext = os.path.splitext(file_path)[1].lower()
    if file_ext == '.ppt':
        # Convert .ppt to .pptx using Aspose.Slides (or any other conversion tool)
        pptx_path = os.path.splitext(file_path)[0] + '.pptx'
        # Conversion code here...
        file_path = pptx_path

    # Define output file path
    output_file = os.path.join(os.path.dirname(file_path), 'output.pptx')

    # Add hyperlink boxes and save the modified PPTX file
    add_hyperlink_box(file_path, output_file)

    # Upload the file to ConvertAPI
    with open(output_file, 'rb') as f:
        response = requests.post(
            'https://v2.convertapi.com/upload',
            headers={'Authorization': f'Bearer {'secret_PmILF2jX8vz85S8T'}'},
            files={'file': f}
        )

    if response.status_code == 200:
        convertapi_url = response.json()
        return Response({"file_url": convertapi_url}, status=status.HTTP_200_OK)
    else:
        return Response({"error": "Failed to upload to ConvertAPI"}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(['GET'])
@permission_classes([AllowAny])
def list_files(request):
    files = UploadedFile.objects.all()
    serializer = UploadedFileSerializer(files, many=True)
    return Response(serializer.data, status=status.HTTP_200_OK)

@api_view(['POST'])
@permission_classes([AllowAny])
def upload_file(request, user_id):
    file = request.FILES.get('file')

    if not file:
        return Response({'error': 'File is required'}, status=status.HTTP_400_BAD_REQUEST)

    try:
        user = User.objects.get(id=user_id)
    except User.DoesNotExist:
        return Response({'error': 'User not found'}, status=status.HTTP_404_NOT_FOUND)

    uploaded_file = UploadedFile(file=file, user=user)
    uploaded_file.save()

    serializer = UploadedFileSerializer(uploaded_file)
    return Response(serializer.data, status=status.HTTP_201_CREATED)

@api_view(['POST'])
@permission_classes([AllowAny])
def delete_file(request, user_id, file_id):
    try:
        user = User.objects.get(id=user_id)
    except User.DoesNotExist:
        return Response({"error": "User not found"}, status=status.HTTP_404_NOT_FOUND)

    try:
        file = UploadedFile.objects.get(id=file_id, user=user)
    except UploadedFile.DoesNotExist:
        return Response({"error": "File not found"}, status=status.HTTP_404_NOT_FOUND)

    # Delete the file from storage
    if os.path.exists(file.file.path):
        os.remove(file.file.path)

    # Delete the database record
    file.delete()

    return Response({"message": "File deleted successfully"}, status=status.HTTP_200_OK) 

@api_view(['GET'])
@permission_classes([AllowAny])
def user_files(request, user_id):
    User = get_user_model()
    try:
        user = User.objects.get(id=user_id)
    except User.DoesNotExist:
        return Response({"error": "User not found"}, status=status.HTTP_404_NOT_FOUND)


    files = UploadedFile.objects.filter(user=user)
    serializer = UploadedFileSerializer(files, many=True)
    return Response(serializer.data, status=status.HTTP_200_OK)